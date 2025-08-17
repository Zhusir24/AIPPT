"""
PPT 生成服务类
负责使用 python-pptx 生成 PowerPoint 文件
"""

import os
import json
import time
import random
import asyncio
from pathlib import Path
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from ..core.config import settings
from .image_service import ImageService


class PPTXService:
    """PPT 生成服务类"""
    
    def __init__(self):
        self.template_dir = Path(settings.TEMPLATE_DIR)
        self.output_dir = Path(settings.UPLOAD_DIR)
        self.static_dir = Path(settings.STATIC_DIR)
        self.output_dir.mkdir(exist_ok=True)
        self.static_dir.mkdir(exist_ok=True)
        
        # 图片服务
        self.image_service = ImageService()
        
        # 模版特定配置
        self.template_configs = {
            "商务蓝": {
                'primary': RGBColor(31, 78, 121),
                'secondary': RGBColor(217, 226, 243),
                'accent': RGBColor(70, 130, 180),
                'background': RGBColor(248, 249, 250),
                'category': '商务',
                'layouts': ['title_content', 'title_content_image', 'two_column', 'image_focus']
            },
            "简约白": {
                'primary': RGBColor(46, 46, 46),
                'secondary': RGBColor(166, 166, 166),
                'accent': RGBColor(102, 102, 102),
                'background': RGBColor(255, 255, 255),
                'category': '简约',
                'layouts': ['title_content', 'clean_layout', 'minimal_image', 'text_focus']
            },
            "活力橙": {
                'primary': RGBColor(255, 107, 53),
                'secondary': RGBColor(247, 147, 30),
                'accent': RGBColor(255, 152, 0),
                'background': RGBColor(255, 248, 230),
                'category': '创意',
                'layouts': ['dynamic_layout', 'creative_image', 'split_content', 'visual_impact']
            },
            "科技紫": {
                'primary': RGBColor(74, 21, 75),
                'secondary': RGBColor(114, 9, 183),
                'accent': RGBColor(139, 69, 19),
                'background': RGBColor(250, 245, 255),
                'category': '科技',
                'layouts': ['tech_layout', 'data_visual', 'innovation_focus', 'modern_grid']
            },
            "自然绿": {
                'primary': RGBColor(46, 125, 50),
                'secondary': RGBColor(76, 175, 80),
                'accent': RGBColor(200, 230, 201),
                'background': RGBColor(248, 255, 248),
                'category': '自然',
                'layouts': ['organic_layout', 'nature_image', 'eco_design', 'green_focus']
            }
        }
        
        # 后备配色方案
        self.color_schemes = {
            1: self.template_configs.get("商务蓝", {}),
            2: self.template_configs.get("科技紫", {}),
            3: self.template_configs.get("自然绿", {})
        }
    
    async def create_presentation(
        self,
        outline: str,
        content_data: Dict[str, Any],
        template_info: Dict[str, Any]
    ) -> Dict[str, Any]:
        """创建演示文稿"""
        print("🚀 开始创建演示文稿")
        print(f"📋 大纲长度: {len(outline)} 字符")
        print(f"🎨 模板信息: {template_info}")
        print(f"📊 内容数据: {len(str(content_data))} 字符")
        
        try:
            # 生成文件名
            timestamp = int(time.time())
            filename = f"generated_ppt_{timestamp}.pptx"
            output_path = self.output_dir / filename
            print(f"📁 生成文件路径: {output_path}")
            
            # 创建演示文稿
            print("📄 创建新的演示文稿对象")
            prs = Presentation()
            
            # 设置幻灯片尺寸为16:9
            print("📐 设置幻灯片尺寸为16:9")
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # 获取模版配置
            template_name = template_info.get('name', '商务蓝')
            template_config = self.template_configs.get(template_name, self.template_configs['商务蓝'])
            print(f"🎨 使用模板: {template_name}")
            
            # 从大纲提取标题
            lines = outline.strip().split('\n')
            presentation_title = lines[0].lstrip('# ').strip() if lines else "AI生成的演示文稿"
            print(f"📝 演示文稿标题: {presentation_title}")
            
            # 创建增强的标题页
            print("🏠 创建标题页...")
            await self._create_enhanced_title_slide_with_images(prs, presentation_title, template_info, template_config)
            
            # 处理生成的内容
            slides_data = content_data.get('slides', [])
            print(f"📑 找到 {len(slides_data)} 个内容页")
            
            if slides_data:
                print("📄 创建内容页...")
                for i, slide_data in enumerate(slides_data):
                    print(f"  - 创建第 {i+1} 页: {slide_data.get('title', '无标题')}")
                    await self._create_diverse_content_slide(prs, slide_data, template_config, i)
            else:
                print("📄 从大纲生成基本幻灯片...")
                # 如果没有详细内容，从大纲生成基本幻灯片
                await self._create_slides_from_outline_enhanced(prs, outline, template_config)
            
            # 创建增强的结束页
            print("🔚 创建结束页...")
            await self._create_enhanced_end_slide_with_images(prs, template_config)
            
            # 保存文件
            print(f"💾 保存PPT文件到: {output_path}")
            prs.save(output_path)
            print("✅ PPT文件保存成功!")
            
            # 获取文件信息
            file_size = output_path.stat().st_size
            slide_count = len(prs.slides)
            
            return {
                'file_path': str(output_path),
                'filename': filename,
                'file_size': file_size,
                'slide_count': slide_count,
                'created_at': time.strftime('%Y-%m-%d %H:%M:%S'),
                'download_url': f"/api/v1/files/download/{filename}",
                'template_used': template_name
            }
            
        except Exception as e:
            raise Exception(f"创建演示文稿失败: {str(e)}")
    
    def _create_enhanced_title_slide(self, prs: Presentation, title: str, template_info: Dict[str, Any], colors: Dict[str, RGBColor]):
        """创建增强的标题页"""
        slide_layout = prs.slide_layouts[0]  # 标题布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors['background']
        
        # 添加装饰形状
        shapes = slide.shapes
        
        # 添加顶部装饰条
        top_bar = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.33), Inches(1.5)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = colors['primary']
        top_bar.line.fill.background()
        
        # 设置标题
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
            # 设置标题样式
            title_frame = title_shape.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.font.name = 'Microsoft YaHei'
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = colors['primary']
            title_para.alignment = PP_ALIGN.CENTER
        
        # 设置副标题
        if len(slide.shapes.placeholders) > 1:
            subtitle_shape = slide.shapes.placeholders[1]
            subtitle_shape.text = f"基于 {template_info.get('name', '智能')} 模板生成\n由 AI-PPTX 自动创建"
            # 设置副标题样式
            subtitle_frame = subtitle_shape.text_frame
            for para in subtitle_frame.paragraphs:
                para.font.name = 'Microsoft YaHei'
                para.font.size = Pt(18)
                para.font.color.rgb = colors['secondary']
                para.alignment = PP_ALIGN.CENTER
    
    def _create_enhanced_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], colors: Dict[str, RGBColor]):
        """创建增强的内容页"""
        slide_layout = prs.slide_layouts[1]  # 内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors['background']
        
        # 添加侧边装饰
        shapes = slide.shapes
        side_bar = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.3), Inches(7.5)
        )
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = colors['accent']
        side_bar.line.fill.background()
        
        # 设置标题
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_data.get("title", "")
            # 设置标题样式
            title_frame = title_shape.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.font.name = 'Microsoft YaHei'
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = colors['primary']
        
        # 设置内容
        if len(slide.shapes.placeholders) > 1:
            content_shape = slide.shapes.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            # 添加要点
            content_list = slide_data.get("content", [])
            if isinstance(content_list, list):
                for i, item in enumerate(content_list):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    # 处理字符串或字典格式的内容
                    if isinstance(item, str):
                        p.text = f"• {item}"
                    elif isinstance(item, dict):
                        p.text = f"• {item.get('point', item.get('text', str(item)))}"
                    else:
                        p.text = f"• {str(item)}"
                    
                    # 设置段落样式
                    p.level = 0
                    p.font.name = 'Microsoft YaHei'
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(51, 51, 51)
                    p.space_after = Pt(12)
    
    def _create_slides_from_outline(self, prs: Presentation, outline: str, colors: Dict[str, RGBColor]):
        """从大纲创建幻灯片"""
        lines = outline.strip().split('\n')
        current_section = None
        current_content = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('## '):
                # 如果有当前章节，先创建幻灯片
                if current_section:
                    self._create_enhanced_simple_slide(prs, current_section, current_content, colors)
                
                # 开始新章节
                current_section = line[3:].strip()
                current_content = []
            elif line.startswith('- '):
                # 添加要点
                current_content.append(line[2:].strip())
        
        # 创建最后一个章节的幻灯片
        if current_section:
            self._create_enhanced_simple_slide(prs, current_section, current_content, colors)
    
    def _create_enhanced_simple_slide(self, prs: Presentation, title: str, content_list: List[str], colors: Dict[str, RGBColor]):
        """创建增强的简单内容幻灯片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors['background']
        
        # 添加装饰元素
        shapes = slide.shapes
        
        # 添加标题背景条
        title_bg = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(0.8),
            Inches(12.33), Inches(1.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = colors['secondary']
        title_bg.line.fill.background()
        
        # 设置标题
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
            title_frame = title_shape.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.font.name = 'Microsoft YaHei'
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = colors['primary']
        
        # 设置内容
        if len(slide.shapes.placeholders) > 1 and content_list:
            content_shape = slide.shapes.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            for i, item in enumerate(content_list[:6]):  # 限制最多6个要点
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"• {item}"
                p.level = 0
                p.font.name = 'Microsoft YaHei'
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(51, 51, 51)
                p.space_after = Pt(12)
    
    def _create_enhanced_end_slide(self, prs: Presentation, colors: Dict[str, RGBColor]):
        """创建增强的结束页"""
        slide_layout = prs.slide_layouts[0]  # 标题布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors['background']
        
        # 添加装饰形状
        shapes = slide.shapes
        
        # 添加底部装饰条
        bottom_bar = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6),
            Inches(13.33), Inches(1.5)
        )
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = colors['primary']
        bottom_bar.line.fill.background()
        
        # 添加感谢圆形
        circle = shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(5.5), Inches(2.5),
            Inches(2.33), Inches(2.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors['accent']
        circle.line.color.rgb = colors['primary']
        circle.line.width = Pt(3)
        
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = "谢谢！"
            title_frame = title_shape.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.font.name = 'Microsoft YaHei'
            title_para.font.size = Pt(48)
            title_para.font.bold = True
            title_para.font.color.rgb = colors['primary']
            title_para.alignment = PP_ALIGN.CENTER
        
        if len(slide.shapes.placeholders) > 1:
            subtitle_shape = slide.shapes.placeholders[1]
            subtitle_shape.text = "感谢您的聆听\n期待您的反馈\n\n由 AI-PPTX 自动生成"
            subtitle_frame = subtitle_shape.text_frame
            for para in subtitle_frame.paragraphs:
                para.font.name = 'Microsoft YaHei'
                para.font.size = Pt(20)
                para.font.color.rgb = colors['secondary']
                para.alignment = PP_ALIGN.CENTER
    
    async def _create_enhanced_title_slide_with_images(self, prs: Presentation, title: str, template_info: Dict[str, Any], template_config: Dict[str, Any]):
        """创建带图片的增强标题页"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置背景
        self._set_slide_background(slide, template_config)
        
        # 添加标题
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title
            self._style_title_text(title_shape, template_config, size=48)
        
        # 添加副标题和作者信息
        if len(slide.shapes.placeholders) > 1:
            subtitle_shape = slide.shapes.placeholders[1]
            subtitle_shape.text = f"基于 AI 技术生成\n模版风格：{template_info.get('name', '未知')}"
            self._style_subtitle_text(subtitle_shape, template_config)
        
        # 跳过图片搜索，避免卡死问题
        try:
            print("⚠️ 暂时跳过背景图片添加，确保PPT生成流畅")
            # category = template_config.get('category', '商务')
            # images = await self.image_service.search_images("background", category, 1)
            # if images:
            #     await self._add_background_image(slide, images[0], template_config)
        except Exception as e:
            print(f"添加背景图片失败: {e}")
        
        print("✅ 标题页创建完成（无背景图片模式）")
    
    async def _create_diverse_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], template_config: Dict[str, Any], slide_index: int):
        """创建多样化内容页"""
        # 选择布局类型
        layouts = template_config.get('layouts', ['title_content'])
        layout_type = layouts[slide_index % len(layouts)]
        
        # 根据布局类型创建不同的幻灯片
        if layout_type in ['title_content_image', 'image_focus', 'creative_image']:
            await self._create_image_content_slide(prs, slide_data, template_config, layout_type)
        elif layout_type in ['two_column', 'split_content']:
            await self._create_two_column_slide(prs, slide_data, template_config)
        elif layout_type in ['data_visual', 'modern_grid']:
            await self._create_data_visual_slide(prs, slide_data, template_config)
        else:
            await self._create_standard_content_slide(prs, slide_data, template_config)
    
    async def _create_image_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], template_config: Dict[str, Any], layout_type: str):
        """创建图片内容页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置背景
        self._set_slide_background(slide, template_config)
        
        # 添加标题
        title = slide_data.get("title", "标题")
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._style_title_text(slide.shapes.title, template_config, size=36)
        
        # 跳过图片搜索，直接使用标准布局
        try:
            print("⚠️ 暂时跳过内容页图片搜索，使用标准布局")
            # keywords = self.image_service.get_image_keywords_from_content(title)
            # category = template_config.get('category', '商务')
            # images = await self.image_service.search_images(keywords[0] if keywords else "business", category, 1)
            
            # 直接使用标准布局，无图片模式
            await self._add_main_content(slide, slide_data, template_config)
            print("✅ 内容页创建完成（无图片模式）")
        except Exception as e:
            print(f"创建内容页失败: {e}")
            # 回退到标准布局
            await self._add_main_content(slide, slide_data, template_config)
    
    async def _create_two_column_slide(self, prs: Presentation, slide_data: Dict[str, Any], template_config: Dict[str, Any]):
        """创建双栏布局幻灯片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置背景
        self._set_slide_background(slide, template_config)
        
        # 添加标题
        title = slide_data.get("title", "标题")
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._style_title_text(slide.shapes.title, template_config, size=36)
        
        # 将内容分成两栏
        content_list = slide_data.get("content", [])
        if isinstance(content_list, list) and len(content_list) > 2:
            mid_point = len(content_list) // 2
            left_content = content_list[:mid_point]
            right_content = content_list[mid_point:]
            
            # 左栏
            left_textbox = slide.shapes.add_textbox(
                Inches(0.5), Inches(2), Inches(6), Inches(5)
            )
            self._add_bullet_content(left_textbox, left_content, template_config)
            
            # 右栏
            right_textbox = slide.shapes.add_textbox(
                Inches(7), Inches(2), Inches(6), Inches(5)
            )
            self._add_bullet_content(right_textbox, right_content, template_config)
        else:
            # 内容不足，使用标准布局
            await self._add_main_content(slide, slide_data, template_config)
    
    async def _create_data_visual_slide(self, prs: Presentation, slide_data: Dict[str, Any], template_config: Dict[str, Any]):
        """创建数据可视化幻灯片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置背景
        self._set_slide_background(slide, template_config)
        
        # 添加标题
        title = slide_data.get("title", "标题")
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._style_title_text(slide.shapes.title, template_config, size=36)
        
        # 跳过数据图表图片搜索
        try:
            print("⚠️ 暂时跳过数据图表图片搜索，确保PPT生成流畅")
            # images = await self.image_service.search_images("chart data visualization", template_config.get('category'), 1)
            # if images:
            #     await self._add_chart_placeholder(slide, images[0], template_config)
        except Exception as e:
            print(f"添加数据图表失败: {e}")
        
        print("✅ 数据页创建完成（无图表图片模式）")
        
        # 添加内容
        await self._add_main_content(slide, slide_data, template_config)
    
    async def _create_standard_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], template_config: Dict[str, Any]):
        """创建标准内容页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置背景
        self._set_slide_background(slide, template_config)
        
        # 添加标题
        title = slide_data.get("title", "标题")
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._style_title_text(slide.shapes.title, template_config, size=36)
        
        # 添加内容
        await self._add_main_content(slide, slide_data, template_config)
    
    async def _create_slides_from_outline_enhanced(self, prs: Presentation, outline: str, template_config: Dict[str, Any]):
        """从大纲创建增强幻灯片"""
        lines = outline.strip().split('\n')
        current_slide_data = None
        slide_count = 0
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            if line.startswith('## '):  # 二级标题 = 新幻灯片
                if current_slide_data:
                    await self._create_diverse_content_slide(prs, current_slide_data, template_config, slide_count)
                    slide_count += 1
                
                current_slide_data = {
                    "title": line[3:].strip(),
                    "content": []
                }
            elif line.startswith('- ') or line.startswith('* '):  # 项目符号
                if current_slide_data:
                    current_slide_data["content"].append(line[2:].strip())
        
        # 处理最后一张幻灯片
        if current_slide_data:
            await self._create_diverse_content_slide(prs, current_slide_data, template_config, slide_count)
    
    async def _create_enhanced_end_slide_with_images(self, prs: Presentation, template_config: Dict[str, Any]):
        """创建带图片的增强结束页"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置背景
        self._set_slide_background(slide, template_config)
        
        # 添加感谢标题
        if slide.shapes.title:
            slide.shapes.title.text = "谢谢！"
            self._style_title_text(slide.shapes.title, template_config, size=48)
        
        # 添加结束语
        if len(slide.shapes.placeholders) > 1:
            subtitle_shape = slide.shapes.placeholders[1]
            subtitle_shape.text = "感谢您的聆听\n期待您的反馈\n\n由 AI-PPTX 自动生成"
            self._style_subtitle_text(subtitle_shape, template_config)
        
        # 跳过装饰性图片，避免卡死问题
        try:
            print("⚠️ 暂时跳过结束页图片添加，确保PPT生成流畅")
            # category = template_config.get('category', '商务')
            # images = await self.image_service.search_images("thank you conclusion", category, 1)
            # if images:
            #     await self._add_decorative_image(slide, images[0], template_config)
        except Exception as e:
            print(f"添加结束页图片失败: {e}")
        
        print("✅ 结束页创建完成（无装饰图片模式）")
    
    # 辅助方法
    def _set_slide_background(self, slide, template_config: Dict[str, Any]):
        """设置幻灯片背景"""
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = template_config.get('background', RGBColor(255, 255, 255))
        except Exception as e:
            print(f"设置背景失败: {e}")
    
    def _style_title_text(self, title_shape, template_config: Dict[str, Any], size: int = 36):
        """设置标题样式"""
        if title_shape and title_shape.has_text_frame:
            text_frame = title_shape.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.font.name = 'Microsoft YaHei'
                paragraph.font.size = Pt(size)
                paragraph.font.bold = True
                paragraph.font.color.rgb = template_config.get('primary', RGBColor(0, 0, 0))
                paragraph.alignment = PP_ALIGN.CENTER
    
    def _style_subtitle_text(self, subtitle_shape, template_config: Dict[str, Any]):
        """设置副标题样式"""
        if subtitle_shape and subtitle_shape.has_text_frame:
            text_frame = subtitle_shape.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.font.name = 'Microsoft YaHei'
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = template_config.get('secondary', RGBColor(128, 128, 128))
                paragraph.alignment = PP_ALIGN.CENTER
    
    async def _add_background_image(self, slide, image_data: Dict[str, Any], template_config: Dict[str, Any]):
        """添加背景图片"""
        try:
            # 这里可以实现背景图片逻辑
            # 由于python-pptx的限制，暂时跳过
            pass
        except Exception as e:
            print(f"添加背景图片失败: {e}")
    
    async def _add_large_image(self, slide, image_data: Dict[str, Any], template_config: Dict[str, Any]):
        """添加大图片"""
        try:
            # 下载并处理图片
            local_path = await self.image_service.download_and_process_image(image_data['url'], (600, 400))
            if local_path and Path(local_path).exists():
                slide.shapes.add_picture(str(local_path), Inches(7), Inches(2), Inches(6), Inches(4))
        except Exception as e:
            print(f"添加大图片失败: {e}")
    
    async def _add_side_image(self, slide, image_data: Optional[Dict[str, Any]], template_config: Dict[str, Any]):
        """添加侧边图片"""
        try:
            if image_data:
                local_path = await self.image_service.download_and_process_image(image_data['url'], (400, 300))
                if local_path and Path(local_path).exists():
                    slide.shapes.add_picture(str(local_path), Inches(8.5), Inches(2.5), Inches(4), Inches(3))
        except Exception as e:
            print(f"添加侧边图片失败: {e}")
    
    async def _add_chart_placeholder(self, slide, image_data: Dict[str, Any], template_config: Dict[str, Any]):
        """添加图表占位符"""
        try:
            local_path = await self.image_service.download_and_process_image(image_data['url'], (500, 350))
            if local_path and Path(local_path).exists():
                slide.shapes.add_picture(str(local_path), Inches(7.5), Inches(2), Inches(5), Inches(3.5))
        except Exception as e:
            print(f"添加图表失败: {e}")
    
    async def _add_decorative_image(self, slide, image_data: Dict[str, Any], template_config: Dict[str, Any]):
        """添加装饰图片"""
        try:
            local_path = await self.image_service.download_and_process_image(image_data['url'], (300, 200))
            if local_path and Path(local_path).exists():
                slide.shapes.add_picture(str(local_path), Inches(10), Inches(5), Inches(3), Inches(2))
        except Exception as e:
            print(f"添加装饰图片失败: {e}")
    
    async def _add_compact_content(self, slide, slide_data: Dict[str, Any], template_config: Dict[str, Any]):
        """添加紧凑内容"""
        content_list = slide_data.get("content", [])
        if content_list:
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(6), Inches(4))
            self._add_bullet_content(textbox, content_list[:4], template_config)  # 限制内容数量
    
    async def _add_main_content(self, slide, slide_data: Dict[str, Any], template_config: Dict[str, Any]):
        """添加主要内容"""
        if len(slide.shapes.placeholders) > 1:
            content_shape = slide.shapes.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            content_list = slide_data.get("content", [])
            self._add_bullet_content_to_frame(text_frame, content_list, template_config)
    
    def _add_bullet_content(self, textbox, content_list: List, template_config: Dict[str, Any]):
        """添加项目符号内容到文本框"""
        text_frame = textbox.text_frame
        text_frame.clear()
        self._add_bullet_content_to_frame(text_frame, content_list, template_config)
    
    def _add_bullet_content_to_frame(self, text_frame, content_list: List, template_config: Dict[str, Any]):
        """添加项目符号内容到文本框架"""
        for i, item in enumerate(content_list):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # 处理不同格式的内容
            if isinstance(item, str):
                p.text = f"• {item}"
            elif isinstance(item, dict):
                p.text = f"• {item.get('point', item.get('text', str(item)))}"
            else:
                p.text = f"• {str(item)}"
            
            # 设置样式
            p.level = 0
            p.font.name = 'Microsoft YaHei'
            p.font.size = Pt(18)
            p.font.color.rgb = template_config.get('primary', RGBColor(51, 51, 51))
            p.space_after = Pt(12)
    
    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """获取文件信息"""
        try:
            path = Path(file_path)
            if not path.exists():
                return None
            
            prs = Presentation(file_path)
            return {
                'file_size': path.stat().st_size,
                'slide_count': len(prs.slides),
                'exists': True
            }
        except Exception:
            return None 